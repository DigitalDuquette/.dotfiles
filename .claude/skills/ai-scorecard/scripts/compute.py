#!/usr/bin/env python
"""AI Scorecard step 1: compute KPIs from the workbook and read the deck's
current last slide, printing both as JSON.

Usage: compute.py

Output JSON keys:
  kpis            - fresh six KPIs from the workbook
  unmapped_ai     - AI values on live rows that map to no group (must be empty)
  sections        - narrative candidate rows per slide section
  slide           - month label and KPI text currently on the deck's last slide
                    (for the anomaly check; null values mean layout drift)
"""
import json
import warnings
from collections import Counter

warnings.filterwarnings("ignore")

import openpyxl
from pptx import Presentation

BASE = ("/Users/jjduqu/Library/CloudStorage/OneDrive-PADNOS/"
        "Information Solutions-Artificial Intelligence - Documents/"
        "Artificial Intelligence")
WORKBOOK = f"{BASE}/AI @ PADNOS.xlsx"
DECK = f"{BASE}/scorecard/AI Scorecard.pptx"

GROUPS = {
    "Agent": {"Agent", "AI-RPA", "Dashboard", "RPA"},
    "Assisted": {"MCP", "Skill"},
    "Support": {"Software"},
}
LIVE = {"Production", "Production - Recent"}
BACKLOG = {"Backlog", "Blocked", "Parking Lot", "Scoping"}
SECTIONS = [
    ("recently_completed", {"Production - Recent"}),
    ("current_priorities", {"UAT", "Development", "Scoping"}),
    ("up_next", {"Open"}),
]


def num(v):
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def group_of(ai):
    for g, members in GROUPS.items():
        if ai in members:
            return g
    return None


def status_of(row, idx):
    return (row[idx["Status"]] or "").strip() if isinstance(row[idx["Status"]], str) else ""


def main():
    wb = openpyxl.load_workbook(WORKBOOK, data_only=True)
    ws = wb["rpa_ai"]
    rows = list(ws.iter_rows(values_only=True))
    idx = {n: i for i, n in enumerate(rows[0]) if n}
    data = [r for r in rows[1:] if any(v not in (None, "") for v in r)]

    live = [r for r in data if status_of(r, idx) in LIVE]
    counts = Counter(group_of(r[idx["AI"]]) for r in live)
    back = [r for r in data
            if status_of(r, idx) in BACKLOG or not status_of(r, idx)]

    result = {
        "kpis": {
            "hours": round(sum(num(r[idx["Hours per Month"]]) for r in live)),
            "agents": counts.get("Agent", 0),
            "assisted": counts.get("Assisted", 0),
            "support": counts.get("Support", 0),
            "backlog_agents": len({r[idx["Project Code"]] for r in back}),
            "backlog_hours": round(sum(num(r[idx["Hours per Month"]]) for r in back)),
        },
        "unmapped_ai": sorted({str(r[idx["AI"]]) for r in live
                               if group_of(r[idx["AI"]]) is None}),
        "sections": {},
    }
    for key, statuses in SECTIONS:
        result["sections"][key] = [
            {
                "area": r[idx["Business Area"]] or "",
                "name": r[idx["Project Name"]],
                "description": r[idx["Description"]],
                "status": status_of(r, idx),
                "ai": r[idx["AI"]],
            }
            for r in data if status_of(r, idx) in statuses
        ]

    # current last slide, for the anomaly comparison
    slide = {"month": None, "hours": None, "agents": None,
             "assisted": None, "support": None,
             "backlog_agents": None, "backlog_hours": None}
    prs = Presentation(DECK)
    shapes = {sh.shape_id: sh for sh in prs.slides[-1].shapes}
    try:
        slide["month"] = shapes[3].text_frame.text.strip()
        kpi = shapes[4].text_frame.paragraphs
        slide["hours"] = kpi[0].runs[0].text
        if len(kpi) >= 9:  # split layout (July 2026 onward)
            slide["agents"] = kpi[3].runs[0].text
            slide["assisted"] = kpi[5].runs[0].text
            slide["support"] = kpi[7].runs[0].text
        backlog = shapes[5].text_frame.paragraphs
        slide["backlog_agents"] = backlog[1].runs[0].text.strip()
        slide["backlog_hours"] = backlog[2].runs[1].text.strip()
    except (KeyError, IndexError):
        pass  # nulls signal layout drift; the skill tells Claude to flag it
    result["slide"] = slide
    result["slide_count"] = len(prs.slides)

    print(json.dumps(result, indent=2, default=str))


if __name__ == "__main__":
    main()
