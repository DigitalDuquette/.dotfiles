#!/usr/bin/env python
"""AI Scorecard step 2: build the month slide on a working copy of the deck.

Usage: update_slide.py <spec.json>

Spec keys:
  workdir       - directory for backup + working copy (session scratchpad)
  mode          - "update" (refresh deck's last slide) or "add" (clone a new one)
  month_label   - e.g. "July 2026"
  kpis          - {hours, agents, assisted, support, backlog_agents, backlog_hours}
  sections      - keys recently_completed / current_priorities / up_next, each a
                  list of [area, description] pairs, or null to leave that box
                  untouched. Blank area = bullet is the description alone.

Copies the deck to <workdir>/scorecard-backup.pptx (untouched) and
<workdir>/scorecard-work.pptx (edited). Never writes the real deck; deploy.sh
does that after PowerPoint validation.

Cloning is python-pptx library-native on purpose. NEVER hand-edit package
XML/rels at the zip level: the monthly slides carry notesSlide and tags
relationships, and copying a rels file verbatim shares a notes part between two
slides, which corrupts the package (PowerPoint "repair" then deletes the
slide). This happened 2026-07.
"""
import copy
import json
import shutil
import sys
import warnings

warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

DECK = ("/Users/jjduqu/Library/CloudStorage/OneDrive-PADNOS/"
        "Information Solutions-Artificial Intelligence - Documents/"
        "Artificial Intelligence/scorecard/AI Scorecard.pptx")
NBSP = "\xa0"
SECTION_SHAPES = {"recently_completed": 28, "current_priorities": 29, "up_next": 30}


def clone_last_slide(prs):
    src = prs.slides[-1]
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):  # drop layout placeholders
        shp._element.getparent().remove(shp._element)
    skip = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in src.shapes._spTree:  # copy shapes only (no notes/tags)
        if child.tag not in skip:
            new.shapes._spTree.append(copy.deepcopy(child))
    rid_map = {rId: new.part.relate_to(rel.target_part, rel.reltype)
               for rId, rel in src.part.rels.items() if "image" in rel.reltype}
    for el in new.shapes._spTree.iter():  # remap image refs
        for attr in (qn("r:embed"), qn("r:link")):
            if el.get(attr) in rid_map:
                el.set(attr, rid_map[el.get(attr)])
    # no dangling relationship ids may remain
    valid = set(new.part.rels)
    for el in new.shapes._spTree.iter():
        for name, value in el.attrib.items():
            if name.startswith("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"):
                assert value in valid, f"dangling relationship {value}"
    return new


def ensure_split_kpi_layout(tf):
    """Guarantee the 9-paragraph split layout (hours + Agents/Assisted/Support).

    Converts the pre-July-2026 5-paragraph layout (hours + single agents pair)
    in place the first time it is met.
    """
    if len(tf.paragraphs) >= 9:
        return
    num_tmpl, lbl_tmpl = tf.paragraphs[3], tf.paragraphs[4]
    num_tmpl.runs[0].font.size = Pt(32)
    lbl_tmpl.runs[0].font.size = Pt(20)
    for _ in range(2):  # add Assisted and Support pairs
        tf._txBody.append(copy.deepcopy(num_tmpl._p))
        tf._txBody.append(copy.deepcopy(lbl_tmpl._p))


def set_kpis(shapes, k):
    tf = shapes[4].text_frame
    ensure_split_kpi_layout(tf)
    p = tf.paragraphs
    p[0].runs[0].text = str(k["hours"])
    p[3].runs[0].text = str(k["agents"])
    p[4].runs[0].text = "Agents"
    p[5].runs[0].text = str(k["assisted"])
    p[6].runs[0].text = "Assisted"
    p[7].runs[0].text = str(k["support"])
    p[8].runs[0].text = "Support"
    b = shapes[5].text_frame.paragraphs
    b[1].runs[0].text = f"{k['backlog_agents']}{NBSP}"
    b[2].runs[1].text = f" {k['backlog_hours']}{NBSP}"


def rebuild_section(shape, bullets):
    tf = shape.text_frame
    tmpl = copy.deepcopy(tf.paragraphs[0]._p)  # bold-area run + plain run
    for para in list(tf.paragraphs):
        tf._txBody.remove(para._p)
    for area, desc in bullets:
        tf._txBody.append(copy.deepcopy(tmpl))
        para = tf.paragraphs[-1]
        for r in para.runs[2:]:
            r._r.getparent().remove(r._r)
        para.runs[0].text = area
        para.runs[1].text = f": {desc}" if area else desc


def dump(shapes):
    for sid in (3, 4, 5, 28, 29, 30):
        print(f"--- shape {sid}")
        for pi, para in enumerate(shapes[sid].text_frame.paragraphs):
            runs = [(r.text, r.font.size.pt if r.font.size else None, r.font.bold)
                    for r in para.runs]
            print(f"  p{pi}: {runs}")


def main():
    spec = json.load(open(sys.argv[1]))
    work = f"{spec['workdir']}/scorecard-work.pptx"
    shutil.copy2(DECK, f"{spec['workdir']}/scorecard-backup.pptx")
    shutil.copy2(DECK, work)

    prs = Presentation(work)
    if spec["mode"] == "add":
        clone_last_slide(prs)
    slide = prs.slides[-1]
    shapes = {sh.shape_id: sh for sh in slide.shapes}
    missing = {3, 4, 5} - set(shapes)
    assert not missing, f"layout drifted: shape ids {missing} not found"

    shapes[3].text_frame.paragraphs[0].runs[0].text = spec["month_label"]
    set_kpis(shapes, spec["kpis"])
    for key, sid in SECTION_SHAPES.items():
        bullets = spec["sections"].get(key)
        if bullets is not None:
            rebuild_section(shapes[sid], bullets)

    prs.save(work)

    verify = Presentation(work)
    print(f"slides: {len(verify.slides)}")
    dump({sh.shape_id: sh for sh in verify.slides[-1].shapes})
    print(f"\nworking copy: {work}")


if __name__ == "__main__":
    main()
